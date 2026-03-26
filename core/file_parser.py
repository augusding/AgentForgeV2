"""AgentForge V2 — 文件解析器，从上传文件中提取文本。"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

_MAX_TEXT = 50000


async def extract_text(file_path: str) -> str:
    """从文件提取文本内容。返回提取到的文本，失败返回描述性信息。"""
    path = Path(file_path)
    if not path.is_file():
        return ""

    suffix = path.suffix.lower()
    try:
        if suffix in (".txt", ".md", ".csv", ".log", ".tsv"):
            return path.read_text(encoding="utf-8", errors="replace")[:_MAX_TEXT]
        elif suffix == ".json":
            data = json.loads(path.read_text(encoding="utf-8"))
            return json.dumps(data, ensure_ascii=False, indent=2)[:_MAX_TEXT]
        elif suffix == ".pdf":
            return _extract_pdf(path)
        elif suffix == ".docx":
            return _extract_docx(path)
        elif suffix in (".xlsx", ".xls"):
            return _extract_xlsx(path)
        elif suffix == ".pptx":
            return _extract_pptx(path)
        elif suffix in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            size = path.stat().st_size
            return f"[图片文件: {path.name}, 大小: {size} bytes, 格式: {suffix[1:].upper()}]"
        elif suffix in (".html", ".htm"):
            return _extract_html(path)
        else:
            try:
                text = path.read_text(encoding="utf-8", errors="replace")[:_MAX_TEXT]
                if text and len(text.strip()) > 10:
                    return text
            except Exception:
                pass
            return f"[不支持的文件类型: {suffix}，文件名: {path.name}]"
    except Exception as e:
        logger.warning("文件解析失败 [%s]: %s", file_path, e)
        return f"[文件解析失败: {path.name}, 错误: {e}]"


def _extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:100]:
                t = page.extract_text() or ""
                if t.strip():
                    texts.append(t)
        return "\n\n".join(texts)[:_MAX_TEXT]
    except ImportError:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(str(path))
            texts = [page.extract_text() or "" for page in reader.pages[:100]]
            return "\n\n".join(t for t in texts if t.strip())[:_MAX_TEXT]
        except ImportError:
            return "[需要安装 pdfplumber 或 PyPDF2]"


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:_MAX_TEXT]
    except ImportError:
        return "[需要安装 python-docx]"


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return "[需要安装 openpyxl: pip install openpyxl]"
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames[:5]:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(max_row=100, values_only=True):
            row_str = "\t".join(str(c) if c is not None else "" for c in row)
            if row_str.strip():
                rows.append(row_str)
        if rows:
            parts.append(f"[工作表: {name}] ({len(rows)} 行)\n" + "\n".join(rows))
    wb.close()
    result = "\n\n".join(parts)
    return result[:_MAX_TEXT] if result else "[空 Excel 文件]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[需要安装 python-pptx: pip install python-pptx]"
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides[:30], 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            parts.append(f"[幻灯片 {i}]\n" + "\n".join(texts))
    result = "\n\n".join(parts)
    return result[:_MAX_TEXT] if result else "[空 PPT 文件]"


def _extract_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r'<script[^>]*>.*?</script>', '', raw, flags=re.DOTALL | re.I)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.I)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text[:_MAX_TEXT]
