"""AgentForge V2 — 文档转换工具。支持 PDF/Word/Excel/PPT/CSV/Markdown 互转。"""

from __future__ import annotations

import csv
import io
import json
import logging
import os
from pathlib import Path

from tools.registry import ToolDefinition

logger = logging.getLogger(__name__)

_ROOT = Path(os.environ.get("AGENTFORGE_ROOT", ".")).resolve()


def _safe(raw: str) -> Path:
    p = (_ROOT / raw).resolve()
    if not str(p).startswith(str(_ROOT)):
        raise PermissionError(f"路径越界: {raw}")
    return p


async def _convert_handler(args: dict) -> str:
    source = args.get("source_path", "")
    fmt = args.get("target_format", "").lower()
    out_raw = args.get("output_path", "")
    content = args.get("content", "")

    if not fmt:
        return json.dumps({"error": "请指定 target_format (docx/pdf/xlsx/pptx/csv/md/txt)"}, ensure_ascii=False)

    # 获取源内容
    text = content
    if source and not content:
        try:
            sp = _safe(source)
            if not sp.is_file():
                return json.dumps({"error": f"源文件不存在: {source}"}, ensure_ascii=False)
            from core.file_parser import extract_text
            text = await extract_text(str(sp))
        except Exception as e:
            return json.dumps({"error": f"读取失败: {e}"}, ensure_ascii=False)

    if not text:
        return json.dumps({"error": "没有可转换的内容"}, ensure_ascii=False)

    out_dir = _ROOT / "data" / "outputs"
    out_dir.mkdir(parents=True, exist_ok=True)
    if not out_raw:
        stem = Path(source).stem if source else "converted"
        out_raw = f"data/outputs/{stem}.{fmt}"

    try:
        op = _safe(out_raw)
        op.parent.mkdir(parents=True, exist_ok=True)
        converters = {"docx": _to_docx, "word": _to_docx, "pdf": _to_pdf, "xlsx": _to_xlsx,
                      "excel": _to_xlsx, "pptx": _to_pptx, "ppt": _to_pptx, "csv": _to_csv,
                      "md": _to_md, "markdown": _to_md, "txt": _to_txt, "text": _to_txt}
        fn = converters.get(fmt)
        if not fn:
            return json.dumps({"error": f"不支持: {fmt}。支持 docx/pdf/xlsx/pptx/csv/md/txt"}, ensure_ascii=False)
        return await fn(text, op)
    except Exception as e:
        return json.dumps({"error": f"转换失败: {e}"}, ensure_ascii=False)


def _result(path: Path, fmt: str, source_text: str = "", **extra) -> str:
    preview = source_text[:500] if source_text else ""
    if not preview and fmt in ("md", "txt", "csv"):
        try: preview = path.read_text(encoding="utf-8", errors="replace")[:800]
        except Exception: pass
    try:
        rel_path = str(path.relative_to(_ROOT)).replace("\\", "/")
    except ValueError:
        rel_path = str(path)
    return json.dumps({"status": "converted", "path": rel_path, "filename": path.name,
                        "size": path.stat().st_size, "format": fmt, "preview": preview, **extra}, ensure_ascii=False)


async def _to_docx(text: str, op: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return json.dumps({"error": "需要 pip install python-docx"}, ensure_ascii=False)
    doc = Document()
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith(("- ", "* ")):
            doc.add_paragraph(s[2:], style="List Bullet")
        else:
            doc.add_paragraph(s)
    op = op.with_suffix(".docx")
    doc.save(str(op))
    return _result(op, "docx", source_text=text)


async def _to_pdf(text: str, op: Path) -> str:
    op = op.with_suffix(".pdf")
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        loaded = False
        for fp in ["C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/simhei.ttf",
                    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"]:
            if Path(fp).exists():
                try:
                    pdf.add_font("CJK", "", fp, uni=True)
                    pdf.set_font("CJK", size=10)
                    loaded = True
                    break
                except Exception:
                    continue
        if not loaded:
            pdf.set_font("Helvetica", size=10)
        for line in text.split("\n"):
            pdf.multi_cell(0, 6, line)
        pdf.output(str(op))
        return _result(op, "pdf", source_text=text)
    except ImportError:
        op = op.with_suffix(".txt")
        op.write_text(text, encoding="utf-8")
        return json.dumps({"status": "converted", "path": str(op), "filename": op.name,
            "size": op.stat().st_size, "format": "txt", "note": "PDF 需 pip install fpdf2，已降级为 txt"}, ensure_ascii=False)


async def _to_xlsx(text: str, op: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return json.dumps({"error": "需要 pip install openpyxl"}, ensure_ascii=False)
    op = op.with_suffix(".xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "数据"
    row_num = 0
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if "\t" in s:
            cells = s.split("\t")
        elif "|" in s and s.count("|") >= 2:
            cells = [c.strip() for c in s.split("|") if c.strip()]
            if all(set(c) <= {"-", " ", ":"} for c in cells):
                continue
        elif "," in s and s.count(",") >= 2:
            cells = s.split(",")
        else:
            cells = [s]
        row_num += 1
        for j, c in enumerate(cells, 1):
            ws.cell(row=row_num, column=j, value=c.strip())
    wb.save(str(op))
    return _result(op, "xlsx", source_text=text)


async def _to_pptx(text: str, op: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return json.dumps({"error": "需要 pip install python-pptx"}, ensure_ascii=False)
    op = op.with_suffix(".pptx")
    prs = Presentation()
    slides: list[dict] = []
    cur: dict = {"title": "", "content": []}
    for line in text.split("\n"):
        s = line.strip()
        if s.startswith("# ") and not s.startswith("## "):
            if cur["title"] or cur["content"]:
                slides.append(cur)
            cur = {"title": s[2:], "content": []}
        elif s.startswith("## "):
            if cur["title"] or cur["content"]:
                slides.append(cur)
            cur = {"title": s[3:], "content": []}
        elif s:
            cur["content"].append(s)
    if cur["title"] or cur["content"]:
        slides.append(cur)
    if not slides:
        slides = [{"title": "内容", "content": text.split("\n")[:30]}]
    for sd in slides[:30]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sd["title"] or "内容"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "\n".join(sd["content"][:20])
    prs.save(str(op))
    return _result(op, "pptx", source_text=text, slides=len(slides))


async def _to_csv(text: str, op: Path) -> str:
    op = op.with_suffix(".csv")
    rows = []
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if "\t" in s:
            rows.append(s.split("\t"))
        elif "|" in s and s.count("|") >= 2:
            cells = [c.strip() for c in s.split("|") if c.strip()]
            if not all(set(c) <= {"-", " ", ":"} for c in cells):
                rows.append(cells)
        elif "," in s:
            rows.append(s.split(","))
        else:
            rows.append([s])
    with open(op, "w", newline="", encoding="utf-8-sig") as f:
        csv.writer(f).writerows(rows)
    return _result(op, "csv", rows=len(rows))


async def _to_md(text: str, op: Path) -> str:
    op = op.with_suffix(".md")
    op.write_text(text, encoding="utf-8")
    return _result(op, "md")


async def _to_txt(text: str, op: Path) -> str:
    op = op.with_suffix(".txt")
    op.write_text(text, encoding="utf-8")
    return _result(op, "txt")


document_converter = ToolDefinition(
    name="document_converter",
    description=(
        "文档格式转换。支持 PDF、Word(docx)、Excel(xlsx)、PPT(pptx)、CSV、Markdown、纯文本之间互转。"
        "当用户要求格式转换（如 PDF转Word、Word转PDF、CSV转Excel）时使用。"
        "通过 source_path 指定源文件或通过 content 传入文本。转换后保存到 data/outputs/。"
    ),
    input_schema={
        "type": "object",
        "properties": {
            "source_path": {"type": "string", "description": "源文件路径（来自附件）。提供 content 时可省略。"},
            "target_format": {"type": "string", "enum": ["docx", "pdf", "xlsx", "pptx", "csv", "md", "txt"],
                              "description": "目标格式"},
            "output_path": {"type": "string", "description": "输出路径（可选，默认 data/outputs/）"},
            "content": {"type": "string", "description": "直接传入文本内容（无 source_path 时使用）"},
        },
        "required": ["target_format"],
    },
    handler=_convert_handler,
    category="document",
)
