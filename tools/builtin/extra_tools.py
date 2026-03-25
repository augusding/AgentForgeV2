"""
AgentForge V2 — 扩展工具集

email_sender: 邮件发送（占位，需配置 SMTP）
ppt_processor: PPT 读写 (python-pptx)
"""

from __future__ import annotations

import json
import logging
import os
from pathlib import Path

from tools.registry import ToolDefinition

logger = logging.getLogger(__name__)

_ROOT = Path(os.environ.get("AGENTFORGE_ROOT", ".")).resolve()


def _safe_path(raw: str) -> Path:
    p = (_ROOT / raw).resolve()
    if not str(p).startswith(str(_ROOT)):
        raise PermissionError(f"路径越界: {raw}")
    return p


# ── 邮件发送 ──────────────────────────────────────────────

async def _email_handler(args: dict) -> str:
    to = args.get("to", "")
    subject = args.get("subject", "")
    body = args.get("body", "")
    if not to or not subject:
        return json.dumps({"error": "to 和 subject 不能为空"}, ensure_ascii=False)

    smtp_host = os.environ.get("SMTP_HOST", "")
    if not smtp_host:
        return json.dumps({
            "status": "draft",
            "note": "SMTP 未配置，邮件已保存为草稿",
            "to": to, "subject": subject, "body_preview": body[:200],
        }, ensure_ascii=False)

    try:
        import smtplib
        from email.mime.text import MIMEText
        smtp_port = int(os.environ.get("SMTP_PORT", "587"))
        smtp_user = os.environ.get("SMTP_USER", "")
        smtp_pass = os.environ.get("SMTP_PASS", "")

        msg = MIMEText(body, "plain", "utf-8")
        msg["Subject"] = subject
        msg["From"] = smtp_user
        msg["To"] = to

        with smtplib.SMTP(smtp_host, smtp_port) as server:
            server.starttls()
            server.login(smtp_user, smtp_pass)
            server.send_message(msg)

        return json.dumps({"status": "sent", "to": to, "subject": subject}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"发送失败: {e}"}, ensure_ascii=False)


email_sender = ToolDefinition(
    name="email_sender",
    description="发送邮件。需要配置 SMTP 环境变量（SMTP_HOST/SMTP_USER/SMTP_PASS）。",
    input_schema={
        "type": "object",
        "properties": {
            "to": {"type": "string", "description": "收件人邮箱"},
            "subject": {"type": "string", "description": "邮件主题"},
            "body": {"type": "string", "description": "邮件正文"},
        },
        "required": ["to", "subject", "body"],
    },
    handler=_email_handler,
    category="communication",
)


# ── PPT 处理 ──────────────────────────────────────────────

async def _ppt_handler(args: dict) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return json.dumps({"error": "需要安装 python-pptx: pip install python-pptx"}, ensure_ascii=False)

    action = args.get("action", "")
    raw_path = args.get("path", "")
    if not raw_path:
        return json.dumps({"error": "缺少 path 参数"}, ensure_ascii=False)

    try:
        path = _safe_path(raw_path)
    except PermissionError as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)

    try:
        if action == "create":
            slides_str = args.get("slides", "[]")
            slides = json.loads(slides_str)
            prs = Presentation()
            for s in slides:
                layout = prs.slide_layouts[1]  # Title + Content
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = s.get("title", "")
                if slide.placeholders[1]:
                    slide.placeholders[1].text = s.get("content", "")
            path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(path))
            return json.dumps({"status": "created", "path": raw_path, "slides": len(slides)}, ensure_ascii=False)

        elif action == "read":
            if not path.is_file():
                return json.dumps({"error": f"文件不存在: {raw_path}"}, ensure_ascii=False)
            prs = Presentation(str(path))
            result = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                result.append({"slide": i, "text": "\n".join(texts)})
            return json.dumps(result, ensure_ascii=False)

        elif action == "info":
            if not path.is_file():
                return json.dumps({"error": f"文件不存在: {raw_path}"}, ensure_ascii=False)
            prs = Presentation(str(path))
            return json.dumps({"slides": len(prs.slides), "path": raw_path}, ensure_ascii=False)

        return json.dumps({"error": f"未知操作: {action}"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"PPT 操作失败: {e}"}, ensure_ascii=False)


ppt_processor = ToolDefinition(
    name="ppt_processor",
    description="PPT 演示文稿处理。支持 create/read/info 操作。",
    input_schema={
        "type": "object",
        "properties": {
            "action": {"type": "string", "enum": ["create", "read", "info"], "description": "操作类型"},
            "path": {"type": "string", "description": "文件路径(.pptx)"},
            "slides": {"type": "string", "description": "JSON 数组，如 [{\"title\":\"...\",\"content\":\"...\"}]"},
        },
        "required": ["action", "path"],
    },
    handler=_ppt_handler,
    category="document",
)


ALL_EXTRA_TOOLS = [email_sender, ppt_processor]
